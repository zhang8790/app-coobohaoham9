# -*- coding: utf-8 -*-
"""
来电有喜 项目 PPT 生成脚本
维度：专业化 · 哲学化 · 落地化
视觉：国潮食养（朱红 + 墨绿 + 宣纸米色）
依赖：python-pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- 调色板 ----------
CREAM  = RGBColor(0xF7, 0xF1, 0xE6)   # 宣纸米背景
RED    = RGBColor(0xD9, 0x48, 0x2B)   # 朱红（喜·温度）
GREEN  = RGBColor(0x2F, 0x5D, 0x3A)   # 墨绿（食养·自然）
INK    = RGBColor(0x2B, 0x26, 0x20)   # 深墨（正文）
GOLD   = RGBColor(0xC9, 0xA2, 0x4B)   # 暖金（点缀）
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
MUTED  = RGBColor(0x7A, 0x6F, 0x60)   # 浅墨（辅助）
LIGHT  = RGBColor(0xEF, 0xE6, 0xD4)   # 浅米卡片
LIGHTG = RGBColor(0xE4, 0xEC, 0xE3)   # 浅绿卡片
DEEP   = RGBColor(0x3A, 0x33, 0x2A)   # 深卡

FONT = "Microsoft YaHei"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def set_ea(run, name=FONT):
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:latin', 'a:ea', 'a:cs'):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set('typeface', name)


def add_text(slide, l, t, w, h, text, size=18, color=INK, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.18):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.04); tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(2)
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color
        set_ea(r)
    return tb


def add_card(slide, l, t, w, h, fill=LIGHT, line=None, radius=0.10):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1.25)
    sp.shadow.inherit = False
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    return sp


def bg(slide, color=CREAM):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def title_bar(slide, title, kicker=None):
    if kicker:
        add_text(slide, 0.72, 0.30, 12, 0.32, kicker, size=13,
                 color=RED, bold=True)
        ty = 0.66
    else:
        ty = 0.52
    add_text(slide, 0.70, ty, 12.0, 0.7, title, size=30, color=INK, bold=True)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72),
                                 Inches(ty + 0.78), Inches(1.15), Inches(0.085))
    bar.fill.solid(); bar.fill.fore_color.rgb = RED; bar.line.fill.background()


def footer(slide, n):
    add_text(slide, 0.70, 7.02, 8, 0.3, "来电有喜  ·  2026 项目全览",
             size=9.5, color=MUTED)
    add_text(slide, 11.8, 7.02, 1.0, 0.3, f"{n:02d}", size=9.5,
             color=MUTED, align=PP_ALIGN.RIGHT)


def bullets(slide, l, t, w, h, items, size=15, color=INK, gap=6, lead_color=RED):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    for i, it in enumerate(items):
        if isinstance(it, tuple):
            txt, lvl = it
        else:
            txt, lvl = it, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.15
        p.space_after = Pt(gap)
        mark = "▪ " if lvl == 0 else "– "
        r0 = p.add_run(); r0.text = mark
        r0.font.size = Pt(size); r0.font.bold = True
        r0.font.color.rgb = lead_color if lvl == 0 else MUTED
        set_ea(r0)
        r = p.add_run(); r.text = txt
        r.font.size = Pt(size); r.font.color.rgb = color
        set_ea(r)
    return tb


def chip(slide, l, t, w, h, text, fill, txt_color=WHITE, size=12, bold=True):
    sp = add_card(slide, l, t, w, h, fill=fill, radius=0.5)
    tf = sp.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = txt_color
    set_ea(r)
    return sp


# =====================================================================
# 1 · 封面
# =====================================================================
s = prs.slides.add_slide(BLANK)
bg(s, CREAM)
deco = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                          Inches(0.32), Inches(7.5))
deco.fill.solid(); deco.fill.fore_color.rgb = RED; deco.line.fill.background()
seal = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(1.7),
                          Inches(1.9), Inches(1.9))
seal.fill.solid(); seal.fill.fore_color.rgb = RED
seal.line.color.rgb = GOLD; seal.line.width = Pt(3)
seal.shadow.inherit = False
tf = seal.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "喜"; r.font.size = Pt(72); r.font.bold = True
r.font.color.rgb = WHITE; set_ea(r)

add_text(s, 1.05, 2.15, 9.2, 1.3, "来电有喜", size=58, color=RED, bold=True)
add_text(s, 1.08, 3.45, 9.4, 0.7,
         "让每一次「吃」都安心 · 合身 · 有温度", size=26, color=INK, bold=True)
add_text(s, 1.08, 4.45, 9.2, 0.6,
         "情绪驱动  ·  本地食养  ·  配料安全扫描", size=17, color=GREEN, bold=True)
add_text(s, 1.08, 5.05, 9.2, 0.5,
         "微信小程序电商平台  |  Taro + React + Supabase", size=14, color=MUTED)
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.08), Inches(6.35),
                        Inches(9.0), Inches(0.02))
ln.fill.solid(); ln.fill.fore_color.rgb = GOLD; ln.line.fill.background()
add_text(s, 1.08, 6.5, 9.2, 0.5, "2026 项目全览  ·  张林", size=13, color=MUTED)

# =====================================================================
# 2 · 哲学序言
# =====================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
title_bar(s, "食物，从来不只是食物", kicker="WHY · 我们为什么出发")
add_text(s, 0.72, 1.7, 11.9, 0.7,
         "现代人面对食物，有三重隐性的焦虑。来电有喜相信：好的食物关系，"
         "是被看见、被理解、被守护。", size=16, color=INK)
cols = [
    ("看不懂", "配料表如天书，添加剂代号背后藏着什么，普通人无从判断。", RED),
    ("吃不对", "同样的食品，对别人是补，对自己可能是负担——体质被长期忽视。", GREEN),
    ("不安心", "食品安全事件频发，信任在一次次犹豫中消耗殆尽。", GOLD),
]
x = 0.72; w = 3.78; gap = 0.27
for i, (h, body, c) in enumerate(cols):
    cx = x + i * (w + gap)
    add_card(s, cx, 2.75, w, 2.5, fill=LIGHT)
    chip(s, cx + 0.25, 3.05, 1.5, 0.55, h, c, size=18)
    add_text(s, cx + 0.28, 3.8, w - 0.55, 1.3, body, size=14.5, color=INK)
add_card(s, 0.72, 5.5, 11.9, 1.15, fill=DEEP)
add_text(s, 1.0, 5.62, 11.4, 0.95,
         "我们的信念：技术不该制造更多焦虑，而应把本属于每个人的「食养常识」"
         "重新还到手里——让选择，重新变得简单而笃定。",
         size=16, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
footer(s, 2)

# =====================================================================
# 3 · 项目定位
# =====================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
title_bar(s, "来电有喜：把食养常识，还给普通人", kicker="WHAT · 我们是什么")
add_text(s, 0.72, 1.65, 11.9, 0.7,
         "一个用技术把「配料安全 · 个性食养 · 情绪陪伴」融为一体的微信小程序——"
         "不卖焦虑，只做你身边的食养守护者。", size=16, color=INK)
pills = [
    ("配料安全扫描", "拍照 / 输入配料表，即刻读懂安全等级与过敏原", RED, LIGHT),
    ("个性化食养推荐", "结合你的健康画像，给出「对您而言」的结论", GREEN, LIGHTG),
    ("情绪匹配导购", "你今天的心情，值得被好好招待", GOLD, LIGHT),
]
x = 0.72; w = 3.78; gap = 0.27
for i, (h, body, c, f) in enumerate(pills):
    cx = x + i * (w + gap)
    add_card(s, cx, 2.7, w, 2.7, fill=f, line=c)
    chip(s, cx + 0.28, 3.0, w - 0.56, 0.6, h, c, size=17)
    add_text(s, cx + 0.3, 3.85, w - 0.6, 1.4, body, size=14.5, color=INK)
footer(s, 3)

# =====================================================================
# 4 · 市场背景（数据）
# =====================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
title_bar(s, "站在三个同时上升的风口上", kicker="MARKET · 时代背景")
data = [
    ("养生零食赛道", "7514", "亿元  ·  2026 年保守预计，逆势增长", RED, 1.0),
    ("药食同源市场", "3700", "亿元  ·  2025 年规模已突破", GREEN, 0.62),
    ("情绪健康食品", "37%", "全球消费者希望零食改善心理健康", GOLD, 0.45),
    ("食品安全 App", "3 亿+", "全球用户，中国占 40% 份额", RED, 0.5),
]
y = 2.0
for name, num, desc, c, frac in data:
    add_card(s, 0.72, y, 11.9, 1.05, fill=LIGHT)
    add_text(s, 0.95, y + 0.12, 3.2, 0.8, name, size=17, color=INK, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 4.2, y + 0.12, 2.6, 0.8, num, size=30, color=c, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, 6.9, y + 0.12, 5.5, 0.8, desc, size=13.5, color=MUTED,
             anchor=MSO_ANCHOR.MIDDLE)
    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.95),
                             Inches(y + 0.86), Inches(11.4 * frac), Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = c; bar.line.fill.background()
    y += 1.18
add_text(s, 0.72, 6.95, 11.9, 0.4,
         "方向是对的：三大趋势叠加，食品电商当然有市场。", size=13,
         color=GREEN, bold=True)
footer(s, 4)

# =====================================================================
# 5 · 用户痛点
# =====================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
title_bar(s, "现代人的三重食品困境", kicker="PAIN · 我们要解决什么")
tri = [
    ("信息差", "看不懂配料", "E 编号、添加剂、营养成分——决策入口就被卡住。", RED),
    ("匹配差", "吃不对体质", "通用推荐千人一面，忽视过敏、慢病与体质差异。", GREEN),
    ("信任差", "买不到安心", "配料隐患与夸大宣传，让「放心买」成为奢侈。", GOLD),
]
x = 0.72; w = 3.78; gap = 0.27
for i, (tag, h, body, c) in enumerate(tri):
    cx = x + i * (w + gap)
    add_card(s, cx, 2.0, w, 3.4, fill=LIGHT, line=c)
    chip(s, cx + 0.3, 2.3, 1.4, 0.5, tag, c, size=14)
    add_text(s, cx + 0.32, 3.0, w - 0.62, 0.7, h, size=22, color=INK, bold=True)
    add_text(s, cx + 0.32, 3.85, w - 0.62, 1.4, body, size=14.5, color=INK)
add_text(s, 0.72, 5.7, 11.9, 0.9,
         "这三道差，本质是「人与食物之间的理解断层」。来电有喜要做的，"
         "就是用技术把断层填平。", size=15, color=INK, bold=True)
footer(s, 5)

# =====================================================================
# 6 · 解决方案总览（四层闭环）
# =====================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
title_bar(s, "一个闭环：从「扫描」到「懂你」", kicker="HOW · 解决方案总览")
steps = [
    ("A · 健康画像", "体质 / 过敏 / 慢病 / 目标\n结构化建模", GREEN),
    ("B · 感知引擎", "analyzeForProfile\n安全·过敏·禁忌·适配", RED),
    ("C · 个性推荐", "对您而言的安全好物\n闭环到下单", GOLD),
    ("D · 可视·学习", "画像看板 · 月度报告\n行为反推画像", GREEN),
]
x = 0.72; w = 2.74; gap = 0.27
for i, (h, body, c) in enumerate(steps):
    cx = x + i * (w + gap)
    add_card(s, cx, 2.3, w, 2.6, fill=LIGHT, line=c)
    chip(s, cx + 0.25, 2.55, w - 0.5, 0.6, h, c, size=15)
    add_text(s, cx + 0.28, 3.35, w - 0.55, 1.4, body, size=13.5, color=INK)
    if i < 3:
        ar = s.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(cx + w - 0.02),
                                Inches(3.35), Inches(0.34), Inches(0.5))
        ar.fill.solid(); ar.fill.fore_color.rgb = GOLD; ar.line.fill.background()
        ar.shadow.inherit = False
add_card(s, 0.72, 5.25, 11.9, 1.25, fill=DEEP)
add_text(s, 1.0, 5.37, 11.4, 1.0,
         "闭环的核心思想：每一次扫描都让系统更懂你，每一次推荐都更贴合你的身体——"
         "数据不再是冷记录，而是「被照顾」的证据。",
         size=15.5, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
footer(s, 6)

# =====================================================================
# 7 · 能力一 配料安全扫描
# =====================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
title_bar(s, "配料安全扫描：扫一下，就知道能不能吃", kicker="CAPABILITY 01")
add_text(s, 0.72, 1.65, 11.9, 0.5,
         "「还不懂配料表？扫一下就知道安不安全、该吃什么。」", size=16,
         color=RED, bold=True)
ents = [("文本输入", RED), ("拍照 OCR", GREEN), ("扫条码", GOLD), ("历史反推", GREEN)]
x = 0.72; w = 2.74; gap = 0.27
for i, (t, c) in enumerate(ents):
    cx = x + i * (w + gap)
    chip(s, cx, 2.35, w, 0.6, t, c, size=15)
add_card(s, 0.72, 3.15, 11.9, 2.0, fill=LIGHT, line=GREEN)
add_text(s, 0.95, 3.35, 11.4, 0.5, "综合安全报告 · ComprehensiveSafetyReport",
         size=15, color=GREEN, bold=True)
bullets(s, 0.95, 3.95, 11.4, 1.1, [
    "添加剂安全等级：逐项解析，标明风险层级",
    "过敏原强预警：扫描项 ∩ 你的过敏清单，红色高危提示",
    "营养与适宜人群：钠糖脂肪一目了然，适配体质标签",
    "合规护栏：全量文案经 shieldCopy() 过滤 + 强制免责声明",
], size=14.5, gap=4)
footer(s, 7)

# =====================================================================
# 8 · 能力二 个性化食养推荐
# =====================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
title_bar(s, "个性化食养：同一张表，不同体质不同结论", kicker="CAPABILITY 02")
add_card(s, 0.72, 1.9, 5.7, 4.4, fill=LIGHTG, line=GREEN)
add_text(s, 1.0, 2.1, 5.2, 0.5, "结构化健康画像", size=17, color=GREEN, bold=True)
bullets(s, 1.0, 2.75, 5.2, 3.3, [
    "生命阶段：儿童 / 成人 / 孕哺 / 老年",
    "过敏清单：11 类致敏原多选",
    "慢病人群：高血压 / 高血糖 / 肠胃虚弱…",
    "身体状态：宫寒 / 易上火 / 脾胃虚寒…",
    "健康目标：控糖 / 护胃 / 助眠 / 减脂…",
], size=14.5, gap=7)
add_card(s, 6.75, 1.9, 6.0, 4.4, fill=LIGHT, line=RED)
add_text(s, 7.05, 2.1, 5.5, 0.5, "「对您而言」的引擎", size=17, color=RED, bold=True)
add_text(s, 7.05, 2.75, 5.5, 1.2,
         "输入配料 + 你的画像 → 输出:\n安全等级 · 过敏命中 · 体质禁忌 · "
         "食养适配分 · 一句话点评。", size=14.5, color=INK)
add_text(s, 7.05, 4.1, 5.5, 0.5, "示例", size=13, color=MUTED, bold=True)
add_text(s, 7.05, 4.45, 5.5, 1.6,
         "「您体寒，本品温热、温中散寒，正相宜」\n"
         "「您高血压，钠偏高，建议少食用」\n"
         "——同一包零食，妈妈和女儿看到的是不同结论。",
         size=14, color=INK, bold=True)
footer(s, 8)

# =====================================================================
# 9 · 能力三 情绪匹配导购
# =====================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
title_bar(s, "情绪匹配导购：你今天的心情，值得被招待", kicker="CAPABILITY 03")
add_text(s, 0.72, 1.6, 11.9, 0.6,
         "这是项目最大的差异化资产——国内几乎没有成熟产品形态的「心情→食物」匹配。",
         size=15, color=INK)
flow = [("输入心情", RED), ("识别情绪类别", GREEN), ("匹配情绪食品", GOLD), ("温暖陪伴下单", GREEN)]
x = 0.72; w = 2.74; gap = 0.27
for i, (t, c) in enumerate(flow):
    cx = x + i * (w + gap)
    add_card(s, cx, 2.5, w, 1.2, fill=LIGHT, line=c)
    add_text(s, cx + 0.2, 2.5, w - 0.4, 1.2, t, size=16, color=INK, bold=True,
             anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    if i < 3:
        ar = s.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(cx + w - 0.02),
                                Inches(2.92), Inches(0.32), Inches(0.4))
        ar.fill.solid(); ar.fill.fore_color.rgb = GOLD; ar.line.fill.background()
        ar.shadow.inherit = False
add_card(s, 0.72, 4.1, 11.9, 2.1, fill=DEEP)
add_text(s, 1.0, 4.25, 11.4, 1.85,
         "关键判断（来自市场验证）：\n\n"
         "情绪匹配不该是「先停下来输入心情再购物」的前台摩擦，而是后台推荐的一个维度——"
         "让用户在日常浏览中，自然被「懂他此刻状态」的食物接住。\n\n"
         "概念独特，但路径要顺：把差异化做成「润物细无声」，而非「刻意为之」。",
         size=15, color=WHITE, bold=True)
footer(s, 9)

# =====================================================================
# 10 · 技术架构
# =====================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
title_bar(s, "技术底座：扎实、可扩展、可运维", kicker="TECH · 技术架构")
tech = [
    ("前端", "Taro 跨端（微信小程序 + H5）\nReact 18 + Zustand + TailwindCSS", RED),
    ("后端", "Supabase：PostgreSQL + RLS 行级安全\nEdge Functions 云函数", GREEN),
    ("云函数", "25+ 个 Edge Functions\n支付 · 分佣 · 退款 · OCR · 红包", GOLD),
    ("数据", "150+ 迁移文件（1~134）\n结构演进完整可追溯", RED),
    ("管理端", "Vite + React + React Router\nadmin-web 完整 PC 后台", GREEN),
]
x = 0.72; w = 3.78; gap = 0.27
for i, (h, body, c) in enumerate(tech):
    cx = x + (i % 3) * (w + gap)
    cy = 2.0 + (i // 3) * 2.2
    add_card(s, cx, cy, w, 1.85, fill=LIGHT, line=c)
    chip(s, cx + 0.25, cy + 0.22, 1.6, 0.55, h, c, size=16)
    add_text(s, cx + 0.28, cy + 0.95, w - 0.55, 0.85, body, size=13, color=INK)
add_text(s, 0.72, 6.5, 11.9, 0.5,
         "约 3 周 / 74 次提交的开发密度——架构清晰、产出扎实，是项目最硬的底气。",
         size=13.5, color=GREEN, bold=True)
footer(s, 10)

# =====================================================================
# 11 · 商业模型
# =====================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
title_bar(s, "商业模型：金豆流通 + 两级推荐 + 商家结算", kicker="BUSINESS · 商业模型")
biz = [
    ("金豆货币体系", "1:1 元，消费抵扣 + 确权返利；\n单位即元，杜绝资损。", RED),
    ("两级推广", "好友（一级）/ 粉丝（二级）；\n每单幂等只发一次，平台保底让利 ×10%。", GREEN),
    ("商家结算", "自营 + 入驻；货款由微信直接打款；\n分佣闭环、对账透明。", GOLD),
]
x = 0.72; w = 3.78; gap = 0.27
for i, (h, body, c) in enumerate(biz):
    cx = x + i * (w + gap)
    add_card(s, cx, 2.0, w, 2.7, fill=LIGHT, line=c)
    chip(s, cx + 0.28, 2.3, w - 0.56, 0.6, h, c, size=16)
    add_text(s, cx + 0.3, 3.15, w - 0.6, 1.4, body, size=14, color=INK)
add_card(s, 0.72, 5.05, 11.9, 1.3, fill=DEEP)
add_text(s, 1.0, 5.17, 11.4, 1.05,
         "「边花边赚」：用户在消费中确权，推广在信任中生长，平台在让利中扩张——"
         "一套让三方都正向循环的食养经济体。",
         size=15.5, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
footer(s, 11)

# =====================================================================
# 12 · 合规护栏
# =====================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
title_bar(s, "合规护栏：食养非医疗，是最高优先级", kicker="COMPLIANCE · 合规红线")
comp = [
    ("shieldCopy()", "违禁词黑名单\n（治疗/治愈/降血压…）\n运行时替换 + 审核双保险", RED),
    ("强制免责声明", "每张分析 / 推荐卡必附\n「食养参考，不替代\n专业医疗诊断」", GREEN),
    ("数据合规 PIPL", "行为分析总闸 + 健康数据\n最小化 + 默认不跨店聚合\n可导出 / 可删除", GOLD),
]
x = 0.72; w = 3.78; gap = 0.27
for i, (h, body, c) in enumerate(comp):
    cx = x + i * (w + gap)
    add_card(s, cx, 2.0, w, 3.0, fill=LIGHT, line=c)
    chip(s, cx + 0.28, 2.3, w - 0.56, 0.6, h, c, size=15)
    add_text(s, cx + 0.3, 3.15, w - 0.6, 1.6, body, size=13.5, color=INK)
add_text(s, 0.72, 5.35, 11.9, 0.9,
         "红线之上，才有增长。合规不是成本，而是食养产品最长的护城河。",
         size=15, color=GREEN, bold=True)
footer(s, 12)

# =====================================================================
# 13 · 市场验证与判断
# =====================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
title_bar(s, "市场判断：有市场，但要会做减法", kicker="VALIDATION · 灰度判断")
add_card(s, 0.72, 2.0, 5.7, 3.2, fill=LIGHTG, line=GREEN)
add_text(s, 1.0, 2.2, 5.2, 0.5, "✓ 三个「有市场」", size=17, color=GREEN, bold=True)
bullets(s, 1.0, 2.85, 5.2, 2.2, [
    "赛道对：情绪健康 + 养生零食 + 食安，三趋势叠加",
    "代码强：74 提交 / 150+ 迁移，架构清晰扎实",
    "概念成立：情绪匹配导购国内几乎空白",
], size=14, gap=9, lead_color=GREEN)
add_card(s, 6.75, 2.0, 6.0, 3.2, fill=LIGHT, line=RED)
add_text(s, 7.05, 2.2, 5.5, 0.5, "✗ 三个「难点」", size=17, color=RED, bold=True)
bullets(s, 7.05, 2.85, 5.5, 2.2, [
    "概念差异 ≠ 用户行为差异（前台输入心情有摩擦）",
    "太全能太散：8 大模块，早期推不动",
    "冷启动鸡生蛋：标签 / 数据库 / 入驻皆需投入",
], size=14, gap=9, lead_color=RED)
add_card(s, 0.72, 5.45, 11.9, 1.1, fill=DEEP)
add_text(s, 1.0, 5.55, 11.4, 0.92,
         "建议：MVP 单城市 · 自营精选 · 入口做窄。  "
         "当前形态市场概率 ≈ 40%  →  重构形态 60–70%。",
         size=15, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
footer(s, 13)

# =====================================================================
# 14 · 落地路线
# =====================================================================
s = prs.slides.add_slide(BLANK)
bg(s)
title_bar(s, "落地路线：先打透一个点", kicker="ROADMAP · 落地路线")
road = [
    ("MVP", "1–2 周", "画像表 → 建档表单 → shieldCopy\nfood-scan 画像感知 → 推荐画像化", RED),
    ("V1.5", "中期", "首页 / 详情画像化\n扫描历史落库", GREEN),
    ("V2", "远期", "健康画像看板 + 月度报告\n扫描行为反推画像", GOLD),
]
x = 0.72; w = 3.78; gap = 0.27
for i, (h, when, body, c) in enumerate(road):
    cx = x + i * (w + gap)
    add_card(s, cx, 2.1, w, 3.1, fill=LIGHT, line=c)
    chip(s, cx + 0.28, 2.4, 1.5, 0.6, h, c, size=18)
    chip(s, cx + 2.0, 2.5, 1.3, 0.45, when, INK, size=12)
    add_text(s, cx + 0.3, 3.35, w - 0.6, 1.6, body, size=14, color=INK)
    if i < 2:
        ar = s.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(cx + w - 0.02),
                                Inches(3.4), Inches(0.32), Inches(0.45))
        ar.fill.solid(); ar.fill.fore_color.rgb = GOLD; ar.line.fill.background()
        ar.shadow.inherit = False
add_text(s, 0.72, 5.5, 11.9, 0.9,
         "战略转向：从一个功能点打透，而不是从十个功能点一起打。",
         size=15, color=INK, bold=True)
footer(s, 14)

# =====================================================================
# 15 · 结语 / 价值宣言
# =====================================================================
s = prs.slides.add_slide(BLANK)
bg(s, DEEP)
# 装饰竖条
deco = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                          Inches(0.32), Inches(7.5))
deco.fill.solid(); deco.fill.fore_color.rgb = RED; deco.line.fill.background()
add_text(s, 1.0, 1.4, 11.3, 0.5, "CLOSING · 价值宣言", size=14,
         color=GOLD, bold=True)
add_text(s, 1.0, 2.1, 11.3, 1.6,
         "我们卖的不是商品，\n是一种更安心的生活方式。", size=36,
         color=WHITE, bold=True, line_spacing=1.1)
add_text(s, 1.0, 4.2, 11.3, 1.4,
         "来电有喜，把「吃」这件每天都发生的小事，重新变得可被理解、可被信任、"
         "可被照顾。\n从卖货，到陪你吃好——这是技术该有的温度。",
         size=17, color=RGBColor(0xE8, 0xE0, 0xD0), line_spacing=1.35)
# 行动条
act = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(5.9),
                         Inches(11.3), Inches(0.8))
act.fill.solid(); act.fill.fore_color.rgb = RED; act.line.fill.background()
act.shadow.inherit = False
tf = act.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "下一步 · MVP 单城灰度，用真实扫描与下单验证闭环"
r.font.size = Pt(17); r.font.bold = True; r.font.color.rgb = WHITE; set_ea(r)
add_text(s, 1.0, 6.95, 11.3, 0.4, "来电有喜  ·  2026 项目全览", size=10,
         color=MUTED)

# ---------- 保存 ----------
OUT = r"C:\Users\zhanglin\Desktop\app-coobohaoham9\deliverables\来电有喜_项目全览.pptx"
prs.save(OUT)
print("SAVED:", OUT, "slides:", len(prs.slides._sldIdLst))
