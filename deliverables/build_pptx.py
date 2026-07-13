# -*- coding: utf-8 -*-
"""来电有喜 · 融资路演 PPT 生成器（python-pptx）
品牌视觉：中国水墨 × 武侠朱砂（宣纸底 / 朱砂 / 墨黑 / 印章签名）
输出：deliverables/来电有喜_融资路演.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- 品牌色 ----
PAPER        = RGBColor(0xFF,0xFB,0xF7)
CINNABAR     = RGBColor(0xC2,0x41,0x0C)
CINNABAR_DEEP= RGBColor(0x9A,0x34,0x12)
INK          = RGBColor(0x1C,0x19,0x17)
SECONDARY    = RGBColor(0x78,0x35,0x0F)
MUTED        = RGBColor(0x8A,0x81,0x78)
LINE         = RGBColor(0xE7,0xDD,0xD0)
CINNABAR_SOFT= RGBColor(0xFF,0xED,0xE3)
WHITE        = RGBColor(0xFF,0xFF,0xFF)
CARD         = RGBColor(0xFF,0xFF,0xFF)

HEAD = 'KaiTi'            # 标题：楷体（武侠古意）
BODY = 'Microsoft YaHei'  # 正文：雅黑

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def set_font(run, name=BODY):
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:ea', 'a:cs'):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set('typeface', name)

def bg(slide):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    s.fill.solid(); s.fill.fore_color.rgb = PAPER
    s.line.fill.background(); s.shadow.inherit = False
    slide.shapes._spTree.remove(s._element)
    slide.shapes._spTree.insert(2, s._element)
    return s

def seal(slide, text, left, top, size=0.62):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(size), Inches(size))
    s.fill.solid(); s.fill.fore_color.rgb = CINNABAR
    s.line.fill.background(); s.shadow.inherit = False
    s.rotation = -3
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.02); tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    lines = text.split('\n')
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = ln
        r.font.size = Pt(int(size * 15)); r.font.bold = True; r.font.color.rgb = WHITE
        set_font(r, HEAD)
    return s

def pageno(slide, n):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11.9), Inches(0.55), Inches(0.55), Pt(2))
    bar.fill.solid(); bar.fill.fore_color.rgb = CINNABAR; bar.line.fill.background(); bar.shadow.inherit = False
    tb = slide.shapes.add_textbox(Inches(11.2), Inches(0.42), Inches(0.7), Inches(0.4))
    p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = f"{n:02d} / 12"; r.font.size = Pt(11); r.font.color.rgb = CINNABAR
    set_font(r)

def eyebrow(slide, text, top=0.62, left=0.9):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(8), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = CINNABAR
    set_font(r)
    # 朱砂短横
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top+0.42), Inches(0.5), Pt(2.5))
    ln.fill.solid(); ln.fill.fore_color.rgb = CINNABAR; ln.line.fill.background(); ln.shadow.inherit = False

def title(slide, text, top=1.1, size=32):
    tb = slide.shapes.add_textbox(Inches(0.9), Inches(top), Inches(11.5), Inches(1.1))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = INK
    set_font(r, HEAD)

def card(slide, left, top, w, h, fill=CARD, soft=False):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = CINNABAR_SOFT if soft else CARD
    s.line.color.rgb = LINE; s.line.width = Pt(1)
    s.shadow.inherit = False
    return s

def textbox(slide, left, top, w, h, lines, size=14, color=INK, bold=False, align=PP_ALIGN.LEFT, font=BODY, space=6):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space)
        if isinstance(ln, tuple):
            txt, kw = ln
        else:
            txt, kw = ln, {}
        r = p.add_run(); r.text = txt
        r.font.size = Pt(kw.get('size', size)); r.font.bold = kw.get('bold', bold)
        r.font.color.rgb = kw.get('color', color)
        set_font(r, kw.get('font', font))
    return tb

def newslide():
    s = prs.slides.add_slide(BLANK)
    bg(s)
    return s

def slide_arrow(slide, left, top):
    a = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(left), Inches(top), Inches(0.35), Inches(0.4))
    a.fill.solid(); a.fill.fore_color.rgb = CINNABAR; a.line.fill.background(); a.shadow.inherit = False
    return a

# ============ 1 封面 ============
s = newslide()
seal(s, "来电\n有喜", 0.95, 2.2, size=1.15)
textbox(s, 2.5, 2.35, 9.5, 1.6, ["来电·有喜"], size=72, color=INK, bold=True, font=HEAD, space=0)
textbox(s, 2.55, 4.05, 9.5, 0.6, ["本地生活的情绪经济平台"], size=24, color=INK, bold=True)
# 标签
tag = card(s, 2.55, 4.8, 4.4, 0.55, soft=True)
textbox(s, 2.7, 4.86, 4.2, 0.45, ["融资介绍 · INVESTOR DECK · 2026"], size=13, color=CINNABAR, bold=True)
textbox(s, 0.95, 6.5, 11.5, 0.8,
        [("把「人与消费之间的情绪连接」数字化", {'color': SECONDARY, 'size': 15}),
         ("武侠 × 朱砂 · 情绪编译引擎 V5", {'color': MUTED, 'size': 13})], space=4)

# ============ 2 目录 ============
s = newslide(); pageno(s, 2); eyebrow(s, "CONTENTS"); title(s, "本场，我们讲五件事")
items = [
    ("01", "被忽视的事实", "为什么本地生活需要被重做"),
    ("02", "三大数字化引擎", "编译 · 导购 · 确权"),
    ("03", "增长飞轮", "双货币 · 段位分销"),
    ("04", "差异化与数字化思维", "我们凭什么不同"),
    ("05", "融资需求", "为什么是现在 · 要什么"),
    ("★", "一页纸 BP", "核心信息速览"),
]
x0, y0, w, h, gx, gy = 0.95, 2.0, 5.6, 1.15, 0.5, 0.35
for i, (num, t, d) in enumerate(items):
    col = i % 2; row = i // 2
    lx = x0 + col * (w + gx); ly = y0 + row * (h + gy)
    card(s, lx, ly, w, h)
    textbox(s, lx + 0.25, ly + 0.18, 1.0, 0.9, [num], size=30, color=CINNABAR, bold=True, font=HEAD)
    textbox(s, lx + 1.15, ly + 0.18, w - 1.3, 0.9,
            [(t, {'size': 17, 'bold': True, 'color': INK}), (d, {'size': 12, 'color': MUTED})], space=3)

# ============ 3 问题 ============
s = newslide(); pageno(s, 3); eyebrow(s, "THE BLIND SPOT")
textbox(s, 0.9, 1.7, 11.6, 3.2,
        ["所有本地生活平台，",
         ("都在数字化「货」和「交易」——", {'color': MUTED}),
         ("却没人数字化「情绪」。", {'color': CINNABAR, 'size': 40})], size=40, color=INK, bold=True, font=HEAD, space=10)
textbox(s, 0.95, 4.9, 10.5, 1.4,
        [("但情绪，才是消费的第一性原理：人先有情绪，才有消费。", {'size': 16, 'color': INK}),
         ("这是最后一片还没被规模化的消费决策维度。", {'size': 16, 'color': INK})], space=8)
seal(s, "喜", 11.6, 5.9, size=0.6)

# ============ 4 定位愿景 ============
s = newslide(); pageno(s, 4); eyebrow(s, "POSITIONING"); title(s, "我们是谁")
card(s, 0.95, 2.2, 5.6, 2.5, soft=True)
textbox(s, 1.2, 2.4, 5.1, 0.4, [("一句话定位", {'color': CINNABAR, 'size': 13, 'bold': True})])
textbox(s, 1.2, 2.95, 5.1, 1.6, [("数字化", {'size': 22, 'color': INK, 'bold': True, 'font': HEAD}),
                                   ("人与消费之间的情绪连接", {'size': 22, 'color': CINNABAR, 'bold': True, 'font': HEAD})], space=4)
card(s, 6.8, 2.2, 5.6, 2.5)
textbox(s, 7.05, 2.4, 5.1, 0.4, [("愿景", {'color': CINNABAR, 'size': 13, 'bold': True})])
textbox(s, 7.05, 2.95, 5.1, 1.6, [("让每一次本地消费，", {'size': 18, 'color': INK}),
                                   ("都成为被接住的情绪片刻。", {'size': 18, 'color': CINNABAR, 'bold': True})], space=6)
textbox(s, 0.95, 5.2, 11.0, 1.0,
        [("不是又一个团购平台，而是一套把「情绪价值」做成可标准化、可量化生产的引擎。", {'size': 15, 'color': MUTED})])

# ============ 5 引擎一 情绪编译 ============
s = newslide(); pageno(s, 5); eyebrow(s, "ENGINE 01 · 供给侧壁垒"); title(s, "情绪编译引擎")
textbox(s, 0.95, 1.95, 11, 0.5, [("把「会写情绪文案」的天赋，变成流水线。", {'size': 16, 'color': INK})])
steps = [
    ("STEP 1", "五维打标", "商家给商品打功能/场景/情绪/身份/感官标签，系统智能推荐"),
    ("STEP 2", "一键编译", "三阶段翻译引擎生成情绪叙事：场景问句→状态确认→身份确认"),
    ("STEP 3", "质量评分", "100 分制实时评分，决定进推荐池还是仅店铺展示"),
]
sw = 3.5; gap = 0.45; sx = 0.95; sy = 2.7
for i, (st, h, d) in enumerate(steps):
    lx = sx + i * (sw + gap)
    card(s, lx, sy, sw, 2.3)
    textbox(s, lx + 0.3, sy + 0.25, sw - 0.6, 0.4, [(st, {'color': CINNABAR, 'size': 12, 'bold': True})])
    textbox(s, lx + 0.3, sy + 0.75, sw - 0.6, 0.5, [(h, {'size': 18, 'bold': True, 'color': INK, 'font': HEAD})])
    textbox(s, lx + 0.3, sy + 1.35, sw - 0.6, 1.0, [(d, {'size': 13, 'color': MUTED})], space=4)
    if i < 2:
        a = slide_arrow(s, lx + sw + 0.05, sy + 1.0)
# 分数环
ring = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(sx), Inches(5.4), Inches(0.95), Inches(0.95))
ring.fill.solid(); ring.fill.fore_color.rgb = WHITE; ring.line.color.rgb = CINNABAR; ring.line.width = Pt(3.5)
ring.shadow.inherit = False
tf = ring.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "100"; r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = CINNABAR; set_font(r)
p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run(); r2.text = "分制"; r2.font.size = Pt(9); r2.font.color.rgb = MUTED; set_font(r2)
textbox(s, sx + 1.2, 5.5, 9.5, 1.0,
        [("标签完整度 30    文案合规性 30    场景精准度 20    确权可达性 20",
          {'size': 14, 'color': INK})], space=4)

# ============ 6 引擎二 情绪导购 ============
s = newslide(); pageno(s, 6); eyebrow(s, "ENGINE 02 · 需求侧差异"); title(s, "情绪导购：从「人找货」到「情绪找货」")
textbox(s, 0.95, 1.95, 11, 0.5, [("用户不搜关键词，而是说出感受——系统识别情绪，推荐场景与就近好物。", {'size': 16, 'color': INK})])
moods = [
    ("😴", "耗竭态", "#A8B5C0", "累·透支"),
    ("🌙", "孤独态", "#C0C8D4", "一个人"),
    ("🎉", "表达驱动", "#B59A2E", "想分享"),
    ("🍃", "平稳态", "#8A7D63", "松弛"),
    ("📷", "怀念态", "#B08C4A", "想从前"),
    ("✨", "渴望态", "#5B7AA8", "想要"),
]
mw = 1.85; mg = 0.18; mx = 0.95; my = 2.8
for i, (e, nm, col, bd) in enumerate(moods):
    lx = mx + i * (mw + mg)
    card(s, lx, my, mw, 2.4)
    textbox(s, lx, my + 0.3, mw, 0.8, [(e, {'size': 34})], align=PP_ALIGN.CENTER)
    textbox(s, lx, my + 1.25, mw, 0.5, [(nm, {'size': 15, 'bold': True, 'color': RGBColor.from_string(col.lstrip('#'))})], align=PP_ALIGN.CENTER)
    textbox(s, lx, my + 1.8, mw, 0.4, [(bd, {'size': 12, 'color': MUTED})], align=PP_ALIGN.CENTER)
textbox(s, 0.95, 5.6, 11.5, 0.6,
        [("说感受 → 识别 6 大情绪态 → 场景卡片 + 附近 Feed 流", {'size': 17, 'color': CINNABAR, 'bold': True})], align=PP_ALIGN.CENTER)

# ============ 7 引擎三 消费确权 ============
s = newslide(); pageno(s, 7); eyebrow(s, "ENGINE 03 · 闭环飞轮"); title(s, "消费即确权")
textbox(s, 0.95, 1.95, 11, 0.5, [("消费结束，关系才开始——用户从过客，变成平台共建者。", {'size': 16, 'color': INK})])
nodes = [("支付成功", "交易完成", True), ("情绪确权", "选 1–6 情绪", False),
         ("情绪豆 + 贡献值", "CV 沉淀", False), ("成长反哺", "分红·徽章", False)]
nw = 2.7; ng = 0.35; nx = 0.95; ny = 2.9
for i, (b, sub, hl) in enumerate(nodes):
    lx = nx + i * (nw + ng)
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(lx), Inches(ny), Inches(nw), Inches(1.4))
    shp.fill.solid(); shp.fill.fore_color.rgb = CINNABAR if hl else CARD
    shp.line.color.rgb = (CINNABAR if hl else LINE); shp.line.width = Pt(1); shp.shadow.inherit = False
    tf = shp.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = b; r.font.size = Pt(15); r.font.bold = True
    r.font.color.rgb = WHITE if hl else INK; set_font(r, HEAD)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = sub; r2.font.size = Pt(11)
    r2.font.color.rgb = (RGBColor(0xFF,0xED,0xE3) if hl else MUTED); set_font(r2)
    if i < 3:
        a = s.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(lx + nw + 0.02), Inches(ny + 0.5), Inches(0.3), Inches(0.4))
        a.fill.solid(); a.fill.fore_color.rgb = CINNABAR; a.line.fill.background(); a.shadow.inherit = False
card(s, 0.95, 4.8, 11.4, 1.6)
textbox(s, 1.2, 5.0, 11.0, 1.4,
        [("确权把单次交易变成长期关系资产：", {'size': 14, 'color': INK, 'bold': True}),
         ("用户为积累贡献值而复购、分享（带邀请码锁客）、拉新——留存与裂变内生于产品机制，而非靠外部补贴。",
          {'size': 14, 'color': INK})], space=6)

# ============ 8 增长飞轮 ============
s = newslide(); pageno(s, 8); eyebrow(s, "GROWTH LOOP"); title(s, "增长飞轮")
cols = [
    ("🪙", "双货币模型", "金豆锚定人民币（1:1，真源，可消费/提现）；情绪豆作激励层，不涉资金池风险。", "资产清晰 · 合规可控"),
    ("⚔️", "段位二级分销", "江湖散修→掌门；L1 佣金 15%–28%、L2 6%–16%，段位越高比例越大、永久生效。", "分享即锁客 · 永久绑定"),
    ("🛡️", "完备风控", "活跃门槛、退款回扣、连续未拓新衰减、严禁刷单——不是资金盘，是合规推广分润。", "可持续 · 抗泡沫"),
]
cw = 3.7; cg = 0.35; cx = 0.95; cy = 2.2
for i, (ic, h, d, k) in enumerate(cols):
    lx = cx + i * (cw + cg)
    card(s, lx, cy, cw, 3.3)
    textbox(s, lx + 0.3, cy + 0.3, 1.0, 0.7, [(ic, {'size': 30})])
    textbox(s, lx + 0.3, cy + 1.1, cw - 0.6, 0.6, [(h, {'size': 18, 'bold': True, 'color': INK, 'font': HEAD})])
    textbox(s, lx + 0.3, cy + 1.7, cw - 0.6, 1.3, [(d, {'size': 13, 'color': MUTED})], space=5)
    textbox(s, lx + 0.3, cy + 2.75, cw - 0.6, 0.4, [(k, {'size': 12, 'bold': True, 'color': CINNABAR})])
textbox(s, 0.95, 5.9, 11.5, 0.7,
        [("情绪共鸣自然分享 → 锁客降获客成本 → 确权提 LTV → 段位放大分销 → 更多内容进推荐池 → 体验更好",
          {'size': 15, 'color': CINNABAR, 'bold': True})], align=PP_ALIGN.CENTER)

# ============ 9 差异化对比 ============
s = newslide(); pageno(s, 9); eyebrow(s, "WHY DIFFERENT"); title(s, "我们 vs 传统本地生活")
left = ["心智：去哪吃玩更便宜", "中心：货（SKU、价格）", "发现：关键词搜索 / 低价推荐",
        "用户：过客、价格敏感者", "留存：补贴、红包、折扣", "增长：买量、地推、补贴战"]
right = ["心智：此刻我需要被照顾", "中心：人（情绪、身份）", "发现：表达情绪 → 场景推荐",
         "用户：侠客（段位、贡献、共建）", "留存：确权资产、成长反哺", "增长：情绪共鸣自然裂变"]
# 左卡
card(s, 0.95, 2.2, 5.6, 4.3)
textbox(s, 1.2, 2.35, 5.1, 0.5, [("传统平台（美团 / 点评 / 抖音团购）", {'size': 14, 'bold': True, 'color': MUTED})])
for i, t in enumerate(left):
    textbox(s, 1.2, 3.0 + i * 0.55, 5.1, 0.5, [(t, {'size': 13.5, 'color': INK})])
# 右卡（朱砂描边）
rc = card(s, 6.8, 2.2, 5.6, 4.3, soft=True)
rc.line.color.rgb = CINNABAR; rc.line.width = Pt(1.5)
textbox(s, 7.05, 2.35, 5.1, 0.5, [("来电有喜", {'size': 14, 'bold': True, 'color': CINNABAR})])
for i, t in enumerate(right):
    p = textbox(s, 7.05, 3.0 + i * 0.55, 5.1, 0.5, [(t, {'size': 13.5, 'color': INK, 'bold': True})])

# ============ 10 数字化思维（表） ============
s = newslide(); pageno(s, 10); eyebrow(s, "DIGITAL THINKING"); title(s, "把情绪，变成数据资产")
rows = [
    ("“这商品很有氛围感”", "五维标签打标", "结构化标签库"),
    ("“文案写得好有情绪”", "三阶段编译引擎", "标准化情绪叙事"),
    ("“这条内容质量不错”", "100 分质量评分", "可质控的推荐池准入"),
    ("“我此刻不知道想要啥”", "6 态情绪识别", "情绪→场景→货匹配"),
    ("“我对这单挺满意”", "消费即确权", "情绪豆 + 贡献值 CV"),
    ("“平台跟我有关”", "成长占比 + 分红估算", "用户共建关系量化"),
]
tbl_shape = s.shapes.add_table(len(rows) + 1, 3, Inches(0.95), Inches(2.2), Inches(11.4), Inches(4.0))
table = tbl_shape.table
table.columns[0].width = Inches(4.0); table.columns[1].width = Inches(3.7); table.columns[2].width = Inches(3.7)
hdr = ["主观模糊的东西", "我们的数字化动作", "产出（数据资产）"]
for c, h in enumerate(hdr):
    cell = table.cell(0, c); cell.fill.solid(); cell.fill.fore_color.rgb = CINNABAR
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = cell.text_frame.paragraphs[0]
    r = p.add_run(); r.text = h; r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = WHITE; set_font(r)
for ri, row in enumerate(rows, start=1):
    for c, val in enumerate(row):
        cell = table.cell(ri, c); cell.fill.solid()
        cell.fill.fore_color.rgb = (CINNABAR_SOFT if ri % 2 == 0 else CARD)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        r = p.add_run(); r.text = val
        r.font.size = Pt(13)
        r.font.color.rgb = (CINNABAR if c == 2 else INK)
        r.font.bold = (c == 0 or c == 2)
        set_font(r)
textbox(s, 0.95, 6.4, 11.5, 0.5,
        [("别人数字化「货」和「交易」，我们数字化「人和货之间的情绪连接」。", {'size': 15, 'color': CINNABAR, 'bold': True})])

# ============ 11 融资需求 ============
s = newslide(); pageno(s, 11); eyebrow(s, "THE ASK"); title(s, "为什么是现在")
card(s, 0.95, 2.2, 6.4, 4.0)
textbox(s, 1.2, 2.4, 5.9, 0.5, [("本轮资金用途", {'size': 16, 'bold': True, 'color': INK})])
uses = ["城市复制 — 把单城模型跑成多城网络",
        "引擎算法深化 — 情绪识别与编译的精度与自动化",
        "品牌势能 — 武侠人格化心智的规模化",
        "供给侧拓展 — 商家情绪编译工作台渗透"]
for i, u in enumerate(uses):
    textbox(s, 1.2, 3.1 + i * 0.72, 5.9, 0.6, [(("• " + u), {'size': 14, 'color': INK})])
card(s, 7.6, 2.2, 4.8, 4.0, soft=True)
textbox(s, 7.85, 2.5, 4.3, 0.4, [("本轮融资", {'size': 13, 'color': MUTED, 'bold': True})])
textbox(s, 7.85, 3.0, 4.3, 1.2, [("¥____ 万 / 轮次", {'size': 40, 'color': CINNABAR, 'bold': True, 'font': HEAD})])
textbox(s, 7.85, 4.4, 4.3, 1.5,
        [("里程碑：V5 已上线 · B/C/治理三端完整", {'size': 13, 'color': INK}),
         ("目标：__ 城 · __ 商家 · __ 月活", {'size': 13, 'color': INK})], space=8)
textbox(s, 0.95, 6.4, 11.5, 0.4, [("* 金额与里程碑为占位，按实际数据填充。", {'size': 12, 'color': MUTED})])

# ============ 12 结语 ============
s = newslide(); pageno(s, 12); bg(s)
seal(s, "来电\n有喜", 0.95, 1.6, size=1.0)
textbox(s, 2.6, 1.7, 9.5, 2.4,
        [("别人数字化货与交易，", {'size': 34, 'color': INK, 'bold': True, 'font': HEAD}),
         ("我们数字化人与货之间的情绪连接。", {'size': 34, 'color': CINNABAR, 'bold': True, 'font': HEAD})], space=10)
textbox(s, 0.95, 4.4, 11.5, 1.0, [("谢谢 · 期待共建", {'size': 28, 'color': CINNABAR, 'bold': True, 'font': HEAD})])
textbox(s, 0.95, 5.3, 11.5, 0.8,
        [("联系方式：____________   |   项目：来电有喜   |   微信小程序 · 本地生活情绪经济", {'size': 14, 'color': MUTED})])

out = "C:/Users/zhanglin/Desktop/app-coobohaoham9/deliverables/来电有喜_融资路演.pptx"
prs.save(out)
print("SAVED:", out, "slides:", len(prs.slides._sldIdLst))
